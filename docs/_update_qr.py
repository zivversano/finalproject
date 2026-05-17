#!/usr/bin/env python3
"""Slide 10: replace both QR images with the working local-network QR
and fix the labels (ngrok is down — don't promise it)."""
import os
from pptx import Presentation

DOCS = os.path.dirname(os.path.abspath(__file__))
QR = os.path.join(DOCS, "app_qr.png")
PPTX = os.path.join(DOCS, "Israel_Transit_Project_Presentation.pptx")
URL = "192.168.1.110:5000/gushdan"

p = Presentation(PPTX)
s = list(p.slides)[9].shapes
shapes = list(s)


def set_run(sh, text):
    para = sh.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r.text = ""


def swap_pic(sh):
    L, T, W, H = sh.left, sh.top, sh.width, sh.height
    sh._element.getparent().remove(sh._element)
    s.add_picture(QR, L, T, W, H)


# [7] home QR pic, [8] home URL, [9] remote header, [10] remote QR pic,
# [11] remote URL label, [32] footer feature line
swap_pic(shapes[7])
set_run(shapes[8], URL)
set_run(shapes[9], "🏠 רשת מקומית (WiFi)")
swap_pic(shapes[10])
set_run(shapes[11], URL)
set_run(shapes[32], "נגיש ברשת המקומית — סריקת קוד QR")

p.save(PPTX)
print("QR slide updated ->", URL)
