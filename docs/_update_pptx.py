#!/usr/bin/env python3
"""Update the project presentation: presenter names + Elastic Cloud migration."""
import copy
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

SRC = "Israel_Transit_Project_Presentation.pptx"
p = Presentation(SRC)
S = p.slides

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
def _q(t): return f"{{{A}}}{t}"


def set_multiline(shape, lines):
    """Rewrite paragraph-0 as multiple lines (a:br separated),
    cloning the original first run's rPr so font/colour/RTL survive."""
    para = shape.text_frame.paragraphs[0]._p
    runs = para.findall(_q("r"))
    tmpl = None
    if runs:
        rpr = runs[0].find(_q("rPr"))
        if rpr is not None:
            tmpl = copy.deepcopy(rpr)
    for el in para.findall(_q("r")) + para.findall(_q("br")):
        para.remove(el)
    end = para.find(_q("endParaRPr"))
    new = []
    for i, line in enumerate(lines):
        if i > 0:
            br = etree.Element(_q("br"))
            if tmpl is not None:
                br.append(copy.deepcopy(tmpl))
            new.append(br)
        r = etree.Element(_q("r"))
        if tmpl is not None:
            r.append(copy.deepcopy(tmpl))
        t = etree.SubElement(r, _q("t"))
        t.text = line
        new.append(r)
    for el in new:
        if end is not None:
            end.addprevious(el)
        else:
            para.append(el)


def oneline(shape, text):
    """Set shape to a single line, keeping run-0 formatting."""
    tf = shape.text_frame
    first = True
    for para in tf.paragraphs:
        if first and para.runs:
            para.runs[0].text = text
            for r in para.runs[1:]:
                r.text = ""
            first = False
        else:
            for r in para.runs:
                r.text = ""
    if first and tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text


def repl(shape, old, new):
    """Replace a substring inside whichever run contains it."""
    for para in shape.text_frame.paragraphs:
        for r in para.runs:
            if old in r.text:
                r.text = r.text.replace(old, new)
                return True
    return False


def sh(slide, idx):
    return list(slide.shapes)[idx]


# ---- Slide 1: presenter names (clone the date box for identical styling) ----
s1 = S[0]
date_box = sh(s1, 4)                       # "פרויקט סיום • מאי 2026"
el = copy.deepcopy(date_box._element)
s1.shapes._spTree.append(el)
new_box = s1.shapes[-1]
new_box.top = Emu(5577840)                 # ~6.1" — below the date line
oneline(new_box, "מגישות:  אורלי בר  •  איריס המבר  •  זיו ורסנו")

# ---- Slide 3: architecture labels -> Elastic Cloud ----
s3 = S[2]
repl(sh(s3, 33), "Search Index", "Elastic Cloud ☁️")
repl(sh(s3, 44), "Dashboards", "Cloud Dashboards")

# ---- Slide 6: storage layer — Elasticsearch is now Elastic Cloud ----
s6 = S[5]
oneline(sh(s6, 10), "Elasticsearch · Elastic Cloud ☁️")
set_multiline(sh(s6, 11), [
    "Managed Service (Elastic Cloud)",
    "Region: eu-central-1 (AWS)",
    "",
    "transit-bus-positions: 2.9M",
    "transit-train-positions: 1.9M",
    "transit-traffic: 592K",
    "transit-service-alerts: 1.6M",
    "",
    "סה״כ: 7M+ רשומות · אפס עומס VM",
])

# ---- Slide 13: live numbers now come from Elastic Cloud ----
s13 = S[12]
oneline(sh(s13, 2),  "מספרים מתוך Elastic Cloud (managed)")
oneline(sh(s13, 5),  "2,917,758")          # bus
oneline(sh(s13, 8),  "1,906,678")          # train
oneline(sh(s13, 11), "591,661")            # traffic
oneline(sh(s13, 14), "1,629,934")          # was trip-updates -> alerts
oneline(sh(s13, 15), "התראות שירות")
oneline(sh(s13, 17), "7,046,031")          # was service-alerts -> grand total
oneline(sh(s13, 18), "סה״כ בענן")

# ---- Slide 16: the ES-OOM challenge now has a real resolution ----
s16 = S[15]
oneline(sh(s16, 15),
        "פתרון: הגירה ל-Elastic Cloud — 7M+ רשומות, אפס עומס על ה-VM")

# ---- Slide 14: tech stack line ----
s14 = S[13]
for j, shape in enumerate(s14.shapes):
    if shape.has_text_frame and "Elasticsearch 8.8" in shape.text_frame.text:
        repl(shape, "Elasticsearch 8.8", "Elasticsearch (Elastic Cloud)")
        break

p.save("Israel_Transit_Project_Presentation.pptx")
print("SAVED. slides:", len(p.slides))
