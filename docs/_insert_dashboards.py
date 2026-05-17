#!/usr/bin/env python3
"""Replace slides 11 & 12 dashboard images with the Elastic Cloud
screenshots, and add a 3rd dashboard slide. Run after the 3 PNGs are
saved in docs/: cloud_buses.png, cloud_trains.png, cloud_delays.png
"""
import copy, os, sys
from pptx import Presentation
from pptx.util import Emu

DOCS = os.path.dirname(os.path.abspath(__file__))
imgs = {k: os.path.join(DOCS, f"cloud_{k}.png")
        for k in ("buses", "trains", "delays")}
missing = [v for v in imgs.values() if not os.path.exists(v)]
if missing:
    sys.exit("MISSING images: " + ", ".join(os.path.basename(m) for m in missing))

PPTX = os.path.join(DOCS, "Israel_Transit_Project_Presentation.pptx")
p = Presentation(PPTX)


def set_run(shape, text):
    para = shape.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for r in para.runs[1:]:
            r.text = ""


def swap_picture(slide, img_path):
    """Replace the first PICTURE on the slide with img_path at the
    SAME left/top/width/height."""
    for sh in list(slide.shapes):
        if sh.shape_type == 13:  # PICTURE
            L, T, W, H = sh.left, sh.top, sh.width, sh.height
            sh._element.getparent().remove(sh._element)
            slide.shapes.add_picture(img_path, L, T, W, H)
            return True
    return False


def duplicate_slide(prs, src_idx):
    """Deep-copy a slide (same layout) and append it; return new slide."""
    src = prs.slides[src_idx]
    layout = src.slide_layout
    new = prs.slides.add_slide(layout)
    # remove default placeholders the layout added
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    # copy every shape element from source
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new


# ---- Slide 11 -> Buses ----
s11 = p.slides[10]
set_run(list(s11.shapes)[1], "📊 דשבורד אוטובוסים — Elastic Cloud")
set_run(list(s11.shapes)[2], "Kibana on Elastic Cloud — Buses (operator breakdown)")
swap_picture(s11, imgs["buses"])

# ---- Slide 12 -> Trains ----
s12 = p.slides[11]
set_run(list(s12.shapes)[1], "🚆 דשבורד רכבות — Elastic Cloud")
set_run(list(s12.shapes)[2], "Kibana on Elastic Cloud — Trains (avg velocity, top lines)")
swap_picture(s12, imgs["trains"])
for sh in list(s12.shapes)[5:]:
    if sh.has_text_frame and "*" in sh.text_frame.text:
        set_run(sh, "* נתוני זמן-אמת מ-Elastic Cloud — 1.9M רשומות רכבת")

# ---- New slide after 12 -> Delays & Alerts ----
nd = duplicate_slide(p, 11)               # clone the (already-updated) Trains slide
set_run(list(nd.shapes)[1], "⚠️ דשבורד עיכובים והתראות — Elastic Cloud")
set_run(list(nd.shapes)[2], "Kibana on Elastic Cloud — Delays & Alerts (by hour of day)")
set_run(list(nd.shapes)[3], "13 / 18")
swap_picture(nd, imgs["delays"])
for sh in list(nd.shapes)[5:]:
    if sh.has_text_frame and "*" in sh.text_frame.text:
        set_run(sh, "* התפלגות התראות לפי שעת היום · 1.6M רשומות · 5,574 קווים")
# move the new slide to position 13 (right after slide 12)
xml = p.slides._sldIdLst
ids = list(xml)
ids.insert(12, ids.pop(-1))
for el in ids:
    xml.append(el)

p.save(PPTX)
print("SAVED. slides:", len(p.slides))
